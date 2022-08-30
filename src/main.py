import requests
from bs4 import BeautifulSoup
from pptx import Presentation

pres = Presentation()

# Add input functionality with the amount of links or command line arguments
url = 'https://en.wikipedia.org/wiki/Python_(programming_language)'

request_results = requests.get(url)

# Grab the main image and the first paragraph(s)

wiki_page = BeautifulSoup(request_results.text, "html.parser")

# Title:
# wiki_page.find("h1", {"id": "firstHeading"}).getText()
# Paragraphs:
# for paragraph in wiki_page.select('p'):
#     print(paragraph.getText())

slide1_register = pres.slide_layouts[1]
slide1 = pres.slides.add_slide(slide1_register)
title1 = slide1.shapes.title
subtitle1 = slide1.placeholders[1]
title1.text = wiki_page.find("h 1", {"id": "firstHeading"}).getText()
subtitle1.text = ""
for paragraph in wiki_page.select('p'):
    subtitle1.text += paragraph.getText()

pres.save("Presentation.pptx")